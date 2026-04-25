import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

class DocumentService:
    def extract_text(self, file_path: str) -> str:
        """
        Extracts raw text from a document based on its extension.
        Supports .pdf, .docx, .pptx, .txt
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pdf':
            return self._extract_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return self._extract_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self._extract_pptx(file_path)
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            raise ValueError(f"Unsupported document format: {ext}")

    def _extract_pdf(self, file_path: str) -> str:
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text() + "\n"
        return text

    def _extract_docx(self, file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    def _extract_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

document_service = DocumentService()
